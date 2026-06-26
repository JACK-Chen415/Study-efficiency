from __future__ import annotations

import csv
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "deliverables"
FIG = OUT / "figures"

TITLE = "基于手机轻量感知与自报告数据的大学生学习效率预测系统设计与实现"
CLASS_NAME = "24级物联网1班"
MEMBER_1 = "陈嘉宏（202400502094）"
MEMBER_2 = "蔡天成（202401100012）"
PPTX_NAME = "202400502094+陈嘉宏+24级物联网1班.pptx"

PRIMARY = "264653"
TEAL = "2A9D8F"
YELLOW = "E9C46A"
ORANGE = "E76F51"
BG = "F7FBFB"
SURFACE = "FFFFFF"
INK = "172033"
MUTED = "657287"
LINE = "D9E4EA"
SOFT_TEAL = "E8F6F1"
SOFT_YELLOW = "FFF4D6"
SOFT_ORANGE = "FCEDE9"


def rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def set_bg(slide, color: str = BG) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(color)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int = 16,
    color: str = INK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def add_title(slide, title: str, idx: int | None = None) -> None:
    add_text(slide, title, 0.5, 0.32, 8.15, 0.48, size=25, color=PRIMARY, bold=True)
    if idx is not None:
        add_page_badge(slide, idx)


def add_page_badge(slide, idx: int) -> None:
    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(9.28), Inches(5.1), Inches(0.38), Inches(0.38))
    badge.fill.solid()
    badge.fill.fore_color.rgb = rgb(ORANGE)
    badge.line.color.rgb = rgb(ORANGE)
    tf = badge.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"{idx:02d}"
    run.font.name = "Arial"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = rgb("FFFFFF")


def add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    body: str,
    *,
    fill: str = SURFACE,
    border: str = LINE,
    title_color: str = PRIMARY,
    body_size: int = 12,
) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(border)
    add_text(slide, title, x + 0.16, y + 0.13, w - 0.32, 0.28, size=14, color=title_color, bold=True)
    add_text(slide, body, x + 0.16, y + 0.48, w - 0.32, h - 0.55, size=body_size, color=INK)


def add_pill(slide, x: float, y: float, w: float, text: str, *, fill: str, color: str = PRIMARY) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(fill)
    add_text(slide, text, x, y + 0.06, w, 0.22, size=10, color=color, bold=True, align=PP_ALIGN.CENTER)


def add_bullets(slide, items: list[str], x: float, y: float, w: float, h: float, *, size: int = 13) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for index, item in enumerate(items):
        para = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        para.text = item
        para.level = 0
        para.font.name = "Microsoft YaHei"
        para.font.size = Pt(size)
        para.font.color.rgb = rgb(INK)
        para.space_after = Pt(7)


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float, *, color: str = PRIMARY) -> None:
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(1.8)
    line.line.end_arrowhead = True


def fit_picture(slide, path: Path, x: float, y: float, w: float, h: float) -> None:
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def metric_data() -> tuple[dict, list[tuple[str, float]]]:
    metrics_path = ROOT / "data/processed/mock/model_metrics_mock_20260618.json"
    fi_path = ROOT / "data/processed/mock/feature_importance_mock_20260618.csv"
    metrics = json.loads(metrics_path.read_text(encoding="utf-8")) if metrics_path.exists() else {}
    features: list[tuple[str, float]] = []
    if fi_path.exists():
        with fi_path.open("r", encoding="utf-8", newline="") as file:
            for row in csv.DictReader(file):
                if len(features) >= 5:
                    break
                features.append((row.get("feature", ""), float(row.get("importance", 0))))
    if not features:
        features = [
            ("noise_level", 0.151744),
            ("phone_distraction", 0.147538),
            ("mood_stress", 0.114447),
            ("light_level", 0.113204),
            ("move_count", 0.089503),
        ]
    return metrics, features


def add_flow_node(slide, x: float, y: float, title: str, body: str, fill: str = SURFACE, border: str = TEAL) -> None:
    add_card(slide, x, y, 1.32, 0.88, title, body, fill=fill, border=border, body_size=8)


def build_pptx() -> Path:
    OUT.mkdir(exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    blank = prs.slide_layouts[6]

    # 1 Title
    slide = prs.slides.add_slide(blank)
    set_bg(slide, "F8FAFC")
    left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(2.75), Inches(5.625))
    left.fill.solid()
    left.fill.fore_color.rgb = rgb(PRIMARY)
    left.line.color.rgb = rgb(PRIMARY)
    add_text(slide, "Python 期末课程设计", 0.35, 0.62, 2.0, 0.35, size=14, color="FFFFFF")
    add_text(slide, "学习效率\n预测系统", 0.35, 1.32, 2.0, 1.3, size=29, color="FFFFFF", bold=True)
    add_text(slide, "真实数据质检\nmock 流程验证", 0.35, 4.22, 2.0, 0.55, size=12, color="FFFFFF")
    add_text(slide, TITLE, 3.25, 1.05, 6.15, 1.35, size=29, color=PRIMARY, bold=True)
    add_text(slide, "手机端学习打卡 · 自报告数据 · 随机森林三分类 · 分析看板", 3.28, 2.72, 6.05, 0.35, size=15, color=INK)
    add_text(slide, f"{CLASS_NAME}\n{MEMBER_1}\n{MEMBER_2}", 3.3, 3.62, 5.2, 0.82, size=12, color=MUTED)

    # 2 Background
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "课程设计背景与目标", 2)
    add_card(slide, 0.55, 1.15, 2.65, 2.75, "为什么做", "学习时长不能解释全部效率差异；同样 60 分钟，可能受任务、疲劳、噪声和手机干扰影响。", fill=SURFACE, border=TEAL)
    add_card(slide, 3.65, 1.15, 2.65, 2.75, "怎么做", "用手机浏览器完成打卡、自报告和可选运动聚合特征采集，后端统一保存。", fill=SOFT_TEAL, border=TEAL)
    add_card(slide, 6.75, 1.15, 2.65, 2.75, "交付目标", "完成可运行、可演示、可解释的课程设计系统，不追求论文级高质量实验。", fill=SOFT_YELLOW, border=YELLOW)
    add_pill(slide, 0.82, 4.55, 8.2, "主线：手机轻量感知 → Python 后端 → 数据处理 → 随机森林预测 → 看板展示", fill="EAF3F2")

    # 3 Architecture
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "系统总体架构", 3)
    nodes = [
        (0.62, "手机浏览器", "打卡/自报告\n可选运动采集"),
        (2.42, "Vue 3 前端", "页面流程\n分析看板"),
        (4.22, "FastAPI 后端", "业务校验\n模型服务"),
        (6.02, "数据库", "users\nsessions\nmotion"),
    ]
    for x, title, body in nodes:
        add_flow_node(slide, x, 1.35, title, body, fill=SURFACE)
    for x in [1.94, 3.74, 5.54]:
        add_arrow(slide, x, 1.79, x + 0.35, 1.79, color=TEAL)
    add_flow_node(slide, 2.42, 3.35, "Python 脚本", "导出 CSV\n质量检查", fill=SOFT_YELLOW, border=YELLOW)
    add_flow_node(slide, 4.22, 3.35, "RandomForest", "三分类\n指标/重要性", fill=SOFT_ORANGE, border=ORANGE)
    add_arrow(slide, 6.7, 2.24, 4.9, 3.34, color=PRIMARY)
    add_arrow(slide, 4.22, 3.79, 3.78, 3.79, color=PRIMARY)
    add_arrow(slide, 4.88, 3.35, 4.88, 2.23, color=PRIMARY)
    add_card(slide, 7.55, 1.15, 1.85, 3.25, "设计边界", "simple-login\n不做复杂认证\n运动特征可缺失\n不阻断主流程", fill="FFFFFF", border=LINE, body_size=11)

    # 4 Collection flow
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "手机端采集流程", 4)
    steps = [
        ("登录", "自选昵称"),
        ("开始", "创建会话"),
        ("学习中", "计时/采集"),
        ("结束表单", "填写量表"),
        ("保存", "生成标签"),
        ("看板", "趋势预测"),
    ]
    for i, (title, body) in enumerate(steps):
        x = 0.5 + i * 1.55
        add_flow_node(slide, x, 1.35, f"{i + 1} {title}", body, fill=SURFACE)
        if i < len(steps) - 1:
            add_arrow(slide, x + 1.32, 1.79, x + 1.53, 1.79, color=TEAL)
    add_card(slide, 0.8, 3.15, 8.45, 1.3, "降级策略", "如果 DeviceMotionEvent 不可用、用户拒绝权限、页面进入后台或运动特征上传失败，系统仍允许提交自报告记录；训练导出阶段使用 motion_available 标记缺失机制。", fill=SOFT_TEAL, border=TEAL, body_size=12)

    # 5 Backend
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "后端 API 与数据存储", 5)
    modules = [
        ("用户", "POST /users/simple-login"),
        ("学习记录", "start / end / list / detail"),
        ("运动特征", "upload / get，按 session upsert"),
        ("模型", "train / predict / metrics / importance"),
        ("分析", "overview / trend / factor-analysis"),
    ]
    for i, (title, body) in enumerate(modules):
        add_card(slide, 0.62 + i * 1.82, 1.18, 1.5, 2.35, title, body, fill=SURFACE if i % 2 else SOFT_TEAL, border=TEAL, body_size=9)
    add_card(slide, 0.72, 4.08, 8.55, 0.75, "P0 数据表", "users、study_sessions、motion_features、predictions；删除历史记录时先写入 deleted_study_sessions 备份表。", fill="FFFFFF", border=LINE, body_size=12)

    # 6 Feature engineering
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "特征工程与模型方法", 6)
    pipeline = [
        ("completed 会话", "排除 abandoned"),
        ("导出 CSV", "左连接运动特征"),
        ("质量检查", "样本/缺失/越界"),
        ("特征工程", "One-Hot + 填补"),
        ("随机森林", "低/中/高三分类"),
        ("输出", "指标/混淆矩阵/重要性"),
    ]
    for i, (title, body) in enumerate(pipeline):
        x = 0.45 + i * 1.56
        fill = SOFT_YELLOW if i < 3 else (SOFT_TEAL if i == 3 else SOFT_ORANGE)
        border = YELLOW if i < 3 else (TEAL if i == 3 else ORANGE)
        add_flow_node(slide, x, 1.25, title, body, fill=fill, border=border)
        if i < len(pipeline) - 1:
            add_arrow(slide, x + 1.32, 1.69, x + 1.54, 1.69, color=PRIMARY)
    add_card(slide, 0.8, 3.08, 3.7, 1.25, "标签转换", "efficiency_score：1-2 → low，3 → medium，4-5 → high", fill=SURFACE, border=LINE)
    add_card(slide, 5.1, 3.08, 3.7, 1.25, "缺失处理", "没有运动数据时保留样本，运动字段填 0，并增加 motion_available。", fill=SURFACE, border=LINE)

    # 7 Data boundary
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "数据来源与实验口径", 7)
    add_card(slide, 0.75, 1.2, 3.9, 2.65, "真实数据", "completed = 17\nlow=4, medium=11, high=2\nmotion 缺失 11 条\n用途：真实采集和质量检查", fill=SURFACE, border=TEAL, body_size=12)
    add_card(slide, 5.25, 1.2, 3.9, 2.65, "mock 数据", "total = 120\nlow/medium/high 各 40 条\n用途：训练、预测、看板流程演示\n不代表真实学习规律", fill=SOFT_ORANGE, border=ORANGE, body_size=12)
    add_pill(slide, 0.9, 4.45, 8.1, "结论边界：真实样本少于 30 条，不形成正式模型有效性结论", fill=SOFT_YELLOW, color=PRIMARY)

    # 8 UI
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "系统界面与看板结果", 8)
    collage = FIG / "fig4-ui-collage.png"
    fit_picture(slide, collage, 1.15, 0.92, 7.65, 4.49)
    add_text(slide, "截图自动采集；看板为 mock/demo 流程展示", 1.15, 5.28, 7.65, 0.18, size=8, color=MUTED, align=PP_ALIGN.CENTER)

    # 9 Model result
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "模型指标与特征重要性", 9)
    metrics, features = metric_data()
    rf = metrics.get("random_forest", {})
    dummy = metrics.get("dummy_baseline", {})
    metric_items = [
        ("Accuracy", rf.get("accuracy", 1.0), dummy.get("accuracy", 0.333333)),
        ("Precision", rf.get("precision_macro", 1.0), dummy.get("precision_macro", 0.111111)),
        ("Recall", rf.get("recall_macro", 1.0), dummy.get("recall_macro", 0.333333)),
        ("F1 Macro", rf.get("f1_macro", 1.0), dummy.get("f1_macro", 0.166667)),
    ]
    for i, (name, value, base) in enumerate(metric_items):
        y = 1.22 + i * 0.58
        add_text(slide, name, 0.65, y, 1.05, 0.22, size=11, color=INK, bold=True)
        add_card(slide, 1.78, y - 0.03, 2.85 * float(value), 0.18, "", "", fill=TEAL, border=TEAL)
        add_text(slide, f"RF {float(value):.3f} / Baseline {float(base):.3f}", 4.82, y - 0.02, 2.3, 0.24, size=10, color=MUTED)
    max_imp = max(v for _, v in features)
    for i, (name, value) in enumerate(features):
        y = 1.05 + i * 0.56
        add_text(slide, name, 6.25, y, 1.55, 0.24, size=10, color=INK)
        add_card(slide, 7.78, y, 1.35 * value / max_imp, 0.17, "", "", fill=YELLOW if i % 2 else TEAL, border=YELLOW if i % 2 else TEAL)
        add_text(slide, f"{value:.3f}", 9.18, y - 0.02, 0.5, 0.22, size=9, color=MUTED)
    add_card(slide, 0.65, 4.23, 8.75, 0.65, "口径说明", "这些指标来自 120 条 mock 数据，只证明 Python 训练、预测和看板链路可运行，不证明真实场景泛化能力。", fill=SOFT_ORANGE, border=ORANGE, body_size=11)

    # 10 Demo
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "演示流程与风险处理", 10)
    add_bullets(slide, ["输入昵称进入系统", "点击开始学习并观察计时状态", "结束后填写地点、任务和量表", "查看历史记录、分析看板和预测建议", "说明真实数据与 mock 数据边界"], 0.8, 1.1, 3.55, 3.75, size=13)
    add_card(slide, 5.0, 1.05, 4.1, 0.92, "传感器兼容性", "运动特征可缺失，不阻断学习记录主流程。", fill=SURFACE, border=TEAL, body_size=11)
    add_card(slide, 5.0, 2.18, 4.1, 0.92, "真实样本不足", "真实 completed=17，只做质检，不包装成模型结论。", fill=SOFT_YELLOW, border=YELLOW, body_size=11)
    add_card(slide, 5.0, 3.31, 4.1, 0.92, "演示环境波动", "准备截图、mock 数据和代码包，保证答辩能讲清链路。", fill=SOFT_ORANGE, border=ORANGE, body_size=11)

    # 11 Summary
    slide = prs.slides.add_slide(blank)
    set_bg(slide, "F8FAFC")
    add_title(slide, "总结与成员分工", 11)
    add_card(slide, 0.75, 1.08, 4.05, 2.75, "完成内容", "手机端打卡、自报告采集、FastAPI 后端、训练数据导出、随机森林三分类、预测接口和分析看板。", fill=SURFACE, border=TEAL, body_size=12)
    add_card(slide, 5.2, 1.08, 4.05, 2.75, "成员分工", f"{MEMBER_1}：系统需求、后端接口、数据字典、模型训练脚本、论文结果整理。\n{MEMBER_2}：前端页面、手机运动采集、看板展示、PPT 制作、演示材料整理。", fill=SURFACE, border=ORANGE, body_size=10)
    add_text(slide, "谢谢老师，请批评指正。", 0.75, 4.42, 8.5, 0.5, size=28, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)

    out = OUT / PPTX_NAME
    prs.save(out)
    return out


if __name__ == "__main__":
    output = build_pptx()
    print(output)
