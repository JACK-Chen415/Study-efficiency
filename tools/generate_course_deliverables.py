from __future__ import annotations

import csv
import json
import math
import os
import shutil
import textwrap
import zipfile
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.platypus import (
    Image as RLImage,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "deliverables"
FIG = OUT / "figures"
ASSET = OUT / "assets"

TITLE = "基于手机轻量感知与自报告数据的大学生学习效率预测系统设计与实现"
SUBTITLE = "24级物联网1班 Python 期末课程设计"
PDF_NAME = "提交人学号+提交人姓名+24级物联网1班.pdf"
PPTX_NAME = "提交人学号+提交人姓名+24级物联网1班.pptx"
ZIP_NAME = "提交人学号+提交人姓名+24级物联网1班_代码.zip"

MEMBER_1 = "成员一姓名（学号：成员一学号）"
MEMBER_2 = "成员二姓名（学号：成员二学号）"
SUBMITTER = "提交人姓名（学号：提交人学号）"
CLASS_NAME = "24级物联网1班"

PRIMARY = "#264653"
SECONDARY = "#2A9D8F"
ACCENT = "#E76F51"
LIGHT = "#E9C46A"
BG = "#F7FBFB"
INK = "#1F2933"
MUTED = "#64748B"


def ensure_dirs() -> None:
    OUT.mkdir(exist_ok=True)
    FIG.mkdir(exist_ok=True)
    ASSET.mkdir(exist_ok=True)


def font_path() -> str:
    candidates = [
        r"C:\Windows\Fonts\msyh.ttc",
        r"C:\Windows\Fonts\simhei.ttf",
        "/mnt/c/Windows/Fonts/msyh.ttc",
        "/mnt/c/Windows/Fonts/simhei.ttf",
    ]
    for item in candidates:
        try:
            if Path(item).exists():
                return item
        except OSError:
            continue
    return ""


FONT_FILE = font_path()


def pil_font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = []
    if bold:
        candidates += [r"C:\Windows\Fonts\msyhbd.ttc", "/mnt/c/Windows/Fonts/msyhbd.ttc"]
    if FONT_FILE:
        candidates.append(FONT_FILE)
    for item in candidates:
        try:
            if Path(item).exists():
                return ImageFont.truetype(item, size=size)
        except OSError:
            continue
    return ImageFont.load_default()


def hex_to_rgb(value: str) -> tuple[int, int, int]:
    value = value.lstrip("#")
    return tuple(int(value[i : i + 2], 16) for i in (0, 2, 4))


def blend(hex_color: str, amount: float = 0.85) -> str:
    r, g, b = hex_to_rgb(hex_color)
    nr = int(r + (255 - r) * amount)
    ng = int(g + (255 - g) * amount)
    nb = int(b + (255 - b) * amount)
    return f"#{nr:02X}{ng:02X}{nb:02X}"


def text_width(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.ImageFont) -> int:
    if not text:
        return 0
    box = draw.textbbox((0, 0), text, font=font)
    return box[2] - box[0]


def wrap_text(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.ImageFont, max_width: int) -> list[str]:
    lines: list[str] = []
    current = ""
    for char in text:
        trial = current + char
        if text_width(draw, trial, font) <= max_width or not current:
            current = trial
        else:
            lines.append(current)
            current = char
    if current:
        lines.append(current)
    return lines


def draw_multiline(
    draw: ImageDraw.ImageDraw,
    xy: tuple[int, int],
    text: str,
    font: ImageFont.ImageFont,
    fill: str,
    max_width: int,
    line_gap: int = 8,
    align: str = "left",
) -> int:
    x, y = xy
    line_height = draw.textbbox((0, 0), "测试Ag", font=font)[3] + line_gap
    for line in wrap_text(draw, text, font, max_width):
        dx = 0
        if align == "center":
            dx = (max_width - text_width(draw, line, font)) // 2
        draw.text((x + dx, y), line, font=font, fill=fill)
        y += line_height
    return y


def draw_box(
    draw: ImageDraw.ImageDraw,
    rect: tuple[int, int, int, int],
    title: str,
    body: str,
    fill: str,
    outline: str,
    title_color: str = PRIMARY,
    radius: int = 20,
) -> None:
    x1, y1, x2, y2 = rect
    draw.rounded_rectangle(rect, radius=radius, fill=fill, outline=outline, width=3)
    title_font = pil_font(30, bold=True)
    body_font = pil_font(22)
    draw.text((x1 + 24, y1 + 22), title, font=title_font, fill=title_color)
    draw_multiline(draw, (x1 + 24, y1 + 68), body, body_font, INK, x2 - x1 - 48, line_gap=8)


def draw_arrow(draw: ImageDraw.ImageDraw, start: tuple[int, int], end: tuple[int, int], color: str = SECONDARY) -> None:
    draw.line([start, end], fill=color, width=5)
    angle = math.atan2(end[1] - start[1], end[0] - start[0])
    size = 18
    points = [
        end,
        (
            int(end[0] - size * math.cos(angle - math.pi / 6)),
            int(end[1] - size * math.sin(angle - math.pi / 6)),
        ),
        (
            int(end[0] - size * math.cos(angle + math.pi / 6)),
            int(end[1] - size * math.sin(angle + math.pi / 6)),
        ),
    ]
    draw.polygon(points, fill=color)


def figure_canvas(title: str) -> tuple[Image.Image, ImageDraw.ImageDraw]:
    img = Image.new("RGB", (1800, 1125), BG)
    draw = ImageDraw.Draw(img)
    draw.rectangle((0, 0, 1800, 110), fill=PRIMARY)
    draw.text((70, 30), title, font=pil_font(44, bold=True), fill="white")
    draw.rectangle((70, 105, 1730, 112), fill=LIGHT)
    return img, draw


def save_figure(img: Image.Image, name: str) -> Path:
    path = FIG / name
    img.save(path, "PNG", optimize=True)
    return path


def create_architecture_figure() -> Path:
    img, draw = figure_canvas("图1  系统总体架构")
    boxes = [
        ((95, 235, 395, 430), "手机浏览器", "开始/结束学习\n自报告表单\n可选运动采集", "#FFFFFF", SECONDARY),
        ((520, 235, 820, 430), "Vue 3 前端", "移动端页面\n历史记录\n分析看板", "#FFFFFF", SECONDARY),
        ((945, 235, 1245, 430), "FastAPI 后端", "接口校验\n业务规则\n模型服务", "#FFFFFF", SECONDARY),
        ((1370, 235, 1670, 430), "数据库", "users\nstudy_sessions\nmotion_features\npredictions", "#FFFFFF", SECONDARY),
        ((520, 660, 820, 855), "Python 脚本", "导出CSV\n质量检查\n训练模型", blend(LIGHT, 0.65), "#D9A441"),
        ((945, 660, 1245, 855), "RandomForest", "低/中/高三分类\n指标\n特征重要性", blend(ACCENT, 0.75), ACCENT),
    ]
    for rect, title, body, fill, outline in boxes:
        draw_box(draw, rect, title, body, fill, outline)
    for start, end in [((395, 332), (520, 332)), ((820, 332), (945, 332)), ((1245, 332), (1370, 332)), ((1520, 430), (1120, 660)), ((945, 760), (820, 760)), ((670, 660), (670, 430))]:
        draw_arrow(draw, start, end)
    draw_multiline(
        draw,
        (120, 940),
        "说明：P0 阶段采用 simple-login，不做复杂认证；手机运动特征是辅助输入，权限拒绝或设备不可用时不阻断学习记录主流程。",
        pil_font(24),
        MUTED,
        1540,
        line_gap=10,
    )
    return save_figure(img, "figure_01_architecture.png")


def create_collection_flow_figure() -> Path:
    img, draw = figure_canvas("图2  手机端学习记录与采集流程")
    steps = [
        ("登录", "昵称进入\n去标识化"),
        ("开始学习", "创建会话\n计时开始"),
        ("前台学习", "可选运动采集\nWake Lock提示"),
        ("结束表单", "地点/任务/量表\n效率评分"),
        ("保存记录", "服务端计算\n时长和标签"),
        ("历史/看板", "查看趋势\n预测与建议"),
    ]
    x = 80
    y = 330
    w = 245
    h = 185
    for idx, (title, body) in enumerate(steps):
        rect = (x + idx * 280, y, x + idx * 280 + w, y + h)
        draw_box(draw, rect, f"{idx + 1}. {title}", body, "#FFFFFF", SECONDARY, radius=18)
        if idx < len(steps) - 1:
            draw_arrow(draw, (rect[2] + 10, y + h // 2), (rect[2] + 35, y + h // 2))
    draw.rounded_rectangle((170, 670, 1630, 855), radius=22, fill=blend(ACCENT, 0.82), outline=ACCENT, width=3)
    draw.text((220, 705), "降级策略", font=pil_font(34, bold=True), fill=PRIMARY)
    draw_multiline(
        draw,
        (220, 760),
        "如果 DeviceMotionEvent 不可用、用户拒绝权限、页面切入后台或上传失败，系统仍允许提交自报告记录，并在训练数据中用 motion_available 标记缺失机制。",
        pil_font(25),
        INK,
        1340,
        line_gap=8,
    )
    return save_figure(img, "figure_02_collection_flow.png")


def create_model_pipeline_figure() -> Path:
    img, draw = figure_canvas("图3  数据清洗与模型训练流程")
    top = [
        ("学习记录库", "completed 会话\n排除 abandoned"),
        ("导出训练 CSV", "自报告字段\n运动聚合字段"),
        ("质量检查", "样本数/标签分布\n缺失与异常"),
        ("特征工程", "类别 One-Hot\n缺失运动填0"),
        ("随机森林", "三分类训练\nDummy baseline"),
        ("输出结果", "指标/混淆矩阵\n特征重要性"),
    ]
    x0, y0, w, h, gap = 95, 300, 235, 170, 50
    for idx, (title, body) in enumerate(top):
        rect = (x0 + idx * (w + gap), y0, x0 + idx * (w + gap) + w, y0 + h)
        draw_box(draw, rect, title, body, "#FFFFFF", SECONDARY, radius=14)
        if idx < len(top) - 1:
            draw_arrow(draw, (rect[2] + 5, y0 + h // 2), (rect[2] + gap - 5, y0 + h // 2))
    draw.rounded_rectangle((130, 655, 805, 850), radius=20, fill=blend(LIGHT, 0.64), outline="#D9A441", width=3)
    draw.text((180, 690), "真实数据口径", font=pil_font(34, bold=True), fill=PRIMARY)
    draw_multiline(draw, (180, 748), "真实 completed = 17，少于 30 条最低训练建议，仅用于质检说明，不形成正式有效性结论。", pil_font(25), INK, 560)
    draw.rounded_rectangle((995, 655, 1670, 850), radius=20, fill=blend(SECONDARY, 0.78), outline=SECONDARY, width=3)
    draw.text((1045, 690), "课程演示口径", font=pil_font(34, bold=True), fill=PRIMARY)
    draw_multiline(draw, (1045, 748), "mock = 120，低/中/高各 40 条，仅用于验证训练、预测、看板链路，不写成真实学习规律。", pil_font(25), INK, 560)
    return save_figure(img, "figure_03_model_pipeline.png")


def create_interface_collage_figure() -> Path:
    img, draw = figure_canvas("图4  手机端与分析看板界面拼图")
    panels = [
        ("首页", "昵称：student_a\n开始学习\n历史记录\n分析看板", (120, 210, 480, 900)),
        ("学习中", "00:42:18\n前台采集\n运动状态：可用\n结束学习", (535, 210, 895, 900)),
        ("结束表单", "地点：图书馆\n任务：编程\n疲劳：3\n效率评分：4", (950, 210, 1310, 900)),
        ("分析看板", "总次数 120\nF1-macro 1.000\n特征重要性\n规则建议", (1365, 210, 1725, 900)),
    ]
    for title, body, rect in panels:
        x1, y1, x2, y2 = rect
        draw.rounded_rectangle(rect, radius=42, fill="#111827", outline="#0F172A", width=5)
        inner = (x1 + 24, y1 + 36, x2 - 24, y2 - 36)
        draw.rounded_rectangle(inner, radius=30, fill="#FFFFFF", outline="#E2E8F0", width=2)
        draw.text((inner[0] + 28, inner[1] + 26), title, font=pil_font(34, bold=True), fill=PRIMARY)
        draw.rounded_rectangle((inner[0] + 28, inner[1] + 90, inner[2] - 28, inner[1] + 170), radius=18, fill=blend(SECONDARY, 0.78), outline=SECONDARY, width=2)
        draw.text((inner[0] + 52, inner[1] + 110), "学习效率预测系统", font=pil_font(22, bold=True), fill=PRIMARY)
        y = inner[1] + 205
        for line in body.split("\n"):
            draw.rounded_rectangle((inner[0] + 28, y, inner[2] - 28, y + 70), radius=12, fill="#F8FAFC", outline="#CBD5E1", width=1)
            draw.text((inner[0] + 52, y + 18), line, font=pil_font(23), fill=INK)
            y += 88
        if title == "分析看板":
            for i, v in enumerate([0.95, 0.78, 0.56, 0.42]):
                bx = inner[0] + 55
                by = inner[1] + 580 + i * 33
                draw.rectangle((bx, by, bx + int(210 * v), by + 18), fill=[SECONDARY, ACCENT, LIGHT, PRIMARY][i])
        else:
            draw.rounded_rectangle((inner[0] + 55, inner[3] - 110, inner[2] - 55, inner[3] - 45), radius=28, fill=ACCENT, outline=ACCENT)
            draw.text((inner[0] + 105, inner[3] - 95), "主要操作", font=pil_font(24, bold=True), fill="white")
    draw.text((120, 970), "注：此图为根据已实现页面流程整理的课程报告界面拼图，展示实际功能入口和状态，不包含个人敏感信息。", font=pil_font(24), fill=MUTED)
    return save_figure(img, "figure_04_interface_collage.png")


def read_mock_metrics() -> dict:
    path = ROOT / "data/processed/mock/model_metrics_mock_20260618.json"
    if path.exists():
        return json.loads(path.read_text(encoding="utf-8"))
    return {}


def read_feature_importance() -> list[tuple[str, float]]:
    path = ROOT / "data/processed/mock/feature_importance_mock_20260618.csv"
    values: list[tuple[str, float]] = []
    if path.exists():
        with path.open("r", encoding="utf-8", newline="") as f:
            for row in csv.DictReader(f):
                if len(values) >= 5:
                    break
                values.append((row["feature"], float(row["importance"])))
    return values or [
        ("noise_level", 0.151744),
        ("phone_distraction", 0.147538),
        ("mood_stress", 0.114447),
        ("light_level", 0.113204),
        ("move_count", 0.089503),
    ]


def create_model_results_figure() -> Path:
    img, draw = figure_canvas("图5  课程演示模型指标与特征重要性")
    metrics = read_mock_metrics()
    rf = metrics.get("random_forest", {})
    dummy = metrics.get("dummy_baseline", {})
    metric_items = [
        ("Accuracy", rf.get("accuracy", 1.0), dummy.get("accuracy", 0.333333)),
        ("Precision", rf.get("precision_macro", 1.0), dummy.get("precision_macro", 0.111111)),
        ("Recall", rf.get("recall_macro", 1.0), dummy.get("recall_macro", 0.333333)),
        ("F1-macro", rf.get("f1_macro", 1.0), dummy.get("f1_macro", 0.166667)),
    ]
    draw.text((120, 190), "指标对比（mock 数据，仅验证流程）", font=pil_font(34, bold=True), fill=PRIMARY)
    x0, y0 = 145, 270
    for idx, (name, value, base) in enumerate(metric_items):
        y = y0 + idx * 92
        draw.text((x0, y), name, font=pil_font(26, bold=True), fill=INK)
        draw.rectangle((x0 + 210, y + 8, x0 + 650, y + 38), fill="#E2E8F0")
        draw.rectangle((x0 + 210, y + 8, x0 + 210 + int(440 * value), y + 38), fill=SECONDARY)
        draw.rectangle((x0 + 210, y + 50, x0 + 650, y + 74), fill="#F1F5F9")
        draw.rectangle((x0 + 210, y + 50, x0 + 210 + int(440 * base), y + 74), fill=ACCENT)
        draw.text((x0 + 675, y), f"RF {value:.3f}", font=pil_font(23), fill=PRIMARY)
        draw.text((x0 + 675, y + 45), f"Baseline {base:.3f}", font=pil_font(21), fill=MUTED)
    draw.text((1010, 190), "Top 5 特征重要性", font=pil_font(34, bold=True), fill=PRIMARY)
    feature_values = read_feature_importance()
    max_imp = max(v for _, v in feature_values)
    for idx, (feature, value) in enumerate(feature_values):
        y = 270 + idx * 82
        draw.text((1010, y), feature, font=pil_font(24, bold=True), fill=INK)
        draw.rectangle((1260, y + 8, 1660, y + 42), fill="#E2E8F0")
        draw.rectangle((1260, y + 8, 1260 + int(400 * value / max_imp), y + 42), fill=LIGHT if idx % 2 else SECONDARY)
        draw.text((1675, y + 3), f"{value:.3f}", font=pil_font(22), fill=MUTED)
    draw.rounded_rectangle((145, 810, 1660, 940), radius=20, fill=blend(ACCENT, 0.84), outline=ACCENT, width=3)
    draw_multiline(
        draw,
        (185, 842),
        "口径说明：真实数据仅 17 条，最低训练样本数检查 FAIL；上方指标来自 120 条 mock 数据，只证明 Python 训练、预测、指标读取和看板展示链路可运行。",
        pil_font(26),
        INK,
        1430,
        line_gap=10,
    )
    return save_figure(img, "figure_05_model_results.png")


def generate_figures() -> list[Path]:
    return [
        create_architecture_figure(),
        create_collection_flow_figure(),
        create_model_pipeline_figure(),
        create_interface_collage_figure(),
        create_model_results_figure(),
    ]


def register_pdf_fonts() -> None:
    pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))


def p(text: str, style: ParagraphStyle) -> Paragraph:
    return Paragraph(text, style)


def pdf_styles() -> dict[str, ParagraphStyle]:
    styles = getSampleStyleSheet()
    base = {
        "fontName": "STSong-Light",
        "textColor": colors.HexColor(INK),
        "leading": 17,
        "fontSize": 10.5,
        "spaceAfter": 8,
    }
    return {
        "title": ParagraphStyle("TitleCN", parent=styles["Title"], fontName="STSong-Light", fontSize=22, leading=30, alignment=TA_CENTER, textColor=colors.HexColor(PRIMARY), spaceAfter=16),
        "subtitle": ParagraphStyle("SubtitleCN", parent=styles["Normal"], fontName="STSong-Light", fontSize=11, leading=17, alignment=TA_CENTER, textColor=colors.HexColor(MUTED), spaceAfter=8),
        "h1": ParagraphStyle("H1CN", parent=styles["Heading1"], fontName="STSong-Light", fontSize=16, leading=22, textColor=colors.HexColor(PRIMARY), spaceBefore=18, spaceAfter=10),
        "h2": ParagraphStyle("H2CN", parent=styles["Heading2"], fontName="STSong-Light", fontSize=13, leading=18, textColor=colors.HexColor(INK), spaceBefore=12, spaceAfter=7),
        "body": ParagraphStyle("BodyCN", parent=styles["BodyText"], alignment=TA_JUSTIFY, **base),
        "small": ParagraphStyle("SmallCN", parent=styles["BodyText"], fontName="STSong-Light", fontSize=9, leading=13, textColor=colors.HexColor(MUTED), spaceAfter=6),
        "caption": ParagraphStyle("CaptionCN", parent=styles["BodyText"], fontName="STSong-Light", fontSize=9, leading=13, alignment=TA_CENTER, textColor=colors.HexColor(MUTED), spaceAfter=10),
        "ref": ParagraphStyle("RefCN", parent=styles["BodyText"], fontName="STSong-Light", fontSize=9.5, leading=14, textColor=colors.HexColor(INK), leftIndent=14, firstLineIndent=-14, spaceAfter=5),
        "callout": ParagraphStyle("CalloutCN", parent=styles["BodyText"], fontName="STSong-Light", fontSize=10.5, leading=16, textColor=colors.HexColor(INK), leftIndent=8, rightIndent=8, spaceBefore=6, spaceAfter=8),
    }


def page_header_footer(canvas, doc) -> None:
    canvas.saveState()
    width, height = A4
    canvas.setStrokeColor(colors.HexColor(SECONDARY))
    canvas.setLineWidth(1)
    canvas.line(2.1 * cm, height - 1.45 * cm, width - 2.1 * cm, height - 1.45 * cm)
    canvas.setFont("STSong-Light", 8)
    canvas.setFillColor(colors.HexColor(MUTED))
    canvas.drawString(2.1 * cm, height - 1.15 * cm, "Python期末课程设计")
    canvas.drawRightString(width - 2.1 * cm, height - 1.15 * cm, "大学生学习效率预测系统")
    canvas.setStrokeColor(colors.HexColor("#CBD5E1"))
    canvas.line(2.1 * cm, 1.45 * cm, width - 2.1 * cm, 1.45 * cm)
    canvas.drawString(2.1 * cm, 1.05 * cm, CLASS_NAME)
    canvas.drawRightString(width - 2.1 * cm, 1.05 * cm, f"第 {doc.page} 页")
    canvas.restoreState()


def make_table(data: list[list[str]], widths: list[float], style_sheet: dict[str, ParagraphStyle]) -> Table:
    body = [[p(cell, style_sheet["small"]) for cell in row] for row in data]
    usable = A4[0] - 4.2 * cm
    table = Table(body, colWidths=[usable * w for w in widths], hAlign="LEFT")
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(PRIMARY)),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTNAME", (0, 0), (-1, -1), "STSong-Light"),
                ("GRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#CBD5E1")),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F8FAFC")]),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ("LEFTPADDING", (0, 0), (-1, -1), 6),
                ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                ("TOPPADDING", (0, 0), (-1, -1), 6),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
            ]
        )
    )
    return table


def add_figure(story: list, path: Path, caption: str, styles: dict[str, ParagraphStyle], width_cm: float = 15.6) -> None:
    img = RLImage(str(path), width=width_cm * cm, height=width_cm * cm * 1125 / 1800)
    story.append(img)
    story.append(p(caption, styles["caption"]))


def build_pdf(figures: list[Path]) -> Path:
    register_pdf_fonts()
    styles = pdf_styles()
    out = OUT / PDF_NAME
    doc = SimpleDocTemplate(
        str(out),
        pagesize=A4,
        rightMargin=2.1 * cm,
        leftMargin=2.1 * cm,
        topMargin=2.0 * cm,
        bottomMargin=2.0 * cm,
        title=TITLE,
        author=SUBMITTER,
    )
    story: list = []
    story.append(Spacer(1, 2.4 * cm))
    story.append(p(TITLE, styles["title"]))
    story.append(p(SUBTITLE, styles["subtitle"]))
    story.append(Spacer(1, 0.6 * cm))
    meta = [
        ["成员", f"{MEMBER_1}；{MEMBER_2}"],
        ["班级", CLASS_NAME],
        ["提交说明", "姓名、学号与分工为占位符，提交前请替换为小组真实信息。"],
    ]
    story.append(make_table(meta, [0.18, 0.82], styles))
    story.append(Spacer(1, 1.1 * cm))
    story.append(p("摘要", styles["h1"]))
    story.append(
        p(
            "本课程设计面向大学生日常自习场景，设计并实现了一个基于手机轻量感知与自报告数据的学习效率预测系统。系统以 Vue 3 移动端网页作为采集入口，记录学习起止时间、地点、任务类型、目标清晰度、环境感受、疲劳程度、压力状态、手机干扰和效率评分，并可选采集手机运动聚合特征。后端采用 Python FastAPI 与 SQLAlchemy 完成数据保存、业务校验、模型训练和预测接口，机器学习部分使用 RandomForestClassifier 实现低/中/高三分类演示。当前真实 completed 样本为 17 条，尚不足以形成正式模型有效性结论；因此本文将真实数据用于质量检查说明，将 120 条 mock 数据用于课程演示链路验证。实验表明，系统已经能够完成学习记录采集、训练数据导出、模型训练、预测结果生成和看板展示等核心流程，适合作为 Python 课程设计的完整原型。",
            styles["body"],
        )
    )
    story.append(p("关键词：学习效率；手机轻量感知；自报告数据；FastAPI；RandomForest；Python", styles["small"]))
    story.append(PageBreak())

    story.append(p("1 引言", styles["h1"]))
    for para in [
        "很多学习记录工具只关注“学习了多久”，但同样 60 分钟的学习可能因为任务难度、目标清晰度、疲劳程度、环境噪声和手机干扰而表现出不同效率。对于大学生自习场景，单纯记录时长难以描述学习过程，也难以形成可解释的数据分析结果。",
        "本项目选择手机浏览器作为轻量入口，不依赖额外硬件，在一次学习结束后通过自报告量表记录用户主观状态，并在设备支持时采集手机运动聚合特征。这样的设计既符合物联网专业“感知—传输—处理—应用”的主线，也能用 Python 后端和机器学习流程展示完整课程设计能力。",
        "本课程设计的目标不是构建生产级学习平台，也不是发表高质量论文，而是在有限时间内完成一个可运行、可演示、可解释的系统原型，并在报告中如实说明数据量、传感器兼容性和模型结论边界。",
    ]:
        story.append(p(para, styles["body"]))

    story.append(p("2 实验方法", styles["h1"]))
    story.append(p("2.1 系统架构", styles["h2"]))
    story.append(p("系统采用前后端分离结构。前端由 Vue 3 和 Vite 实现移动端页面，负责 simple-login、开始学习、计时、结束表单、历史记录和分析看板；后端由 FastAPI 提供 REST API，负责用户、学习会话、运动特征、预测结果和分析接口；模型训练脚本读取学习记录与运动特征，生成训练 CSV，完成质量检查、One-Hot 编码和随机森林三分类训练。", styles["body"]))
    add_figure(story, figures[0], "图1 系统总体架构：前端采集、后端存储、Python 特征工程与模型服务组成完整闭环。", styles)

    story.append(p("2.2 数据采集字段", styles["h2"]))
    table1 = [
        ["类别", "字段", "说明"],
        ["学习会话", "start_time、end_time、duration_minutes、time_period", "记录学习起止时间、持续分钟数和开始时间所属时段。"],
        ["自报告字段", "location、task_type、goal_clarity、light_level、noise_level、fatigue_level、mood_stress、phone_distraction、efficiency_score", "描述学习地点、任务、目标、环境、身心状态、手机干扰和主观效率评分。"],
        ["运动特征", "move_count、shake_count、still_ratio、avg_acceleration、max_acceleration", "由 DeviceMotionEvent 聚合得到，仅作为辅助输入，允许缺失。"],
        ["模型标签", "efficiency_label", "由 efficiency_score 派生：1-2 为 low，3 为 medium，4-5 为 high。"],
    ]
    story.append(make_table(table1, [0.16, 0.39, 0.45], styles))
    story.append(p("手机运动数据不是专注度的直接判断依据，只作为辅助特征。权限拒绝、浏览器不支持、页面后台暂停或上传失败时，主流程仍允许提交学习记录。", styles["body"]))
    add_figure(story, figures[1], "图2 手机端学习记录与采集流程：传感器不可用时走降级路径，不阻断自报告数据保存。", styles)

    story.append(p("2.3 后端接口与数据存储", styles["h2"]))
    table2 = [
        ["模块", "核心接口/文件", "功能"],
        ["用户", "POST /api/users/simple-login", "按昵称创建或进入用户，不实现复杂认证。"],
        ["学习记录", "POST /api/sessions/start、POST /api/sessions/end、GET /api/sessions/list", "开始学习、结束学习、查看历史记录。"],
        ["运动特征", "POST /api/motion/upload、GET /api/motion/{session_id}", "上传或读取聚合运动特征。"],
        ["模型", "POST /api/model/train、POST /api/model/predict、GET /api/model/metrics", "训练随机森林、预测效率等级、读取模型指标。"],
        ["分析", "GET /api/analytics/overview、trend、factor-analysis", "提供总览统计、趋势、特征重要性和建议。"],
    ]
    story.append(make_table(table2, [0.16, 0.42, 0.42], styles))

    story.append(p("2.4 特征工程与模型方法", styles["h2"]))
    story.append(p("训练数据由 study_sessions 左连接 motion_features 得到。已放弃会话和无标签记录被排除；缺失运动特征保留样本，运动字段填 0，并增加 motion_available 表示是否采集到运动数据。类别字段 time_period、location 和 task_type 使用 One-Hot 编码，数值字段包括学习时长、六个自报告量表、五个运动字段和 motion_available。主模型采用 RandomForestClassifier，同时使用多数类 DummyClassifier 作为 baseline。", styles["body"]))
    add_figure(story, figures[2], "图3 数据清洗与模型训练流程：真实数据用于质检，mock 数据用于课程演示。", styles)

    story.append(p("3 结果", styles["h1"]))
    story.append(p("3.1 系统功能实现", styles["h2"]))
    story.append(p("前端已经实现 simple-login、首页、学习中计时页、结束学习表单、历史记录、分析看板以及传感器不可用降级提示。后端已经实现学习记录、运动特征、模型训练、预测和分析看板 API，并通过接口测试验证核心业务状态。", styles["body"]))
    add_figure(story, figures[3], "图4 手机端与分析看板界面拼图：展示首页、学习中、结束表单和分析看板主要状态。", styles)

    story.append(p("3.2 数据质量与实验口径", styles["h2"]))
    table3 = [
        ["数据类型", "样本量/分布", "用途", "结论边界"],
        ["真实数据", "completed=17；low=4、medium=11、high=2；motion 缺失 11 条", "用于说明真实采集与质量检查", "少于 30 条，不形成正式模型有效性结论。"],
        ["mock 数据", "120 条；low/medium/high 各 40 条", "用于训练、预测、看板流程演示", "仅证明系统链路可运行，不代表真实学习规律。"],
        ["测试/演示记录", "开发过程产生", "接口与页面验证", "不得计入真实实验结果。"],
    ]
    story.append(make_table(table3, [0.16, 0.28, 0.25, 0.31], styles))
    story.append(p("真实数据检查结果显示异常学习时长、自报告越界、枚举值问题和标签转换不一致均为 0，但总样本数未达到 30 条最低训练建议。因此本文采用课程演示口径，避免将 mock 指标包装成真实结论。", styles["body"]))

    story.append(p("3.3 模型与看板演示结果", styles["h2"]))
    story.append(p("在 120 条 mock 数据上，随机森林训练、指标读取、特征重要性展示和预测 API 链路均可运行。由于 mock 数据是按规则构造的，Accuracy 和 F1-macro 达到 1.000 只说明流程验证通过，不能证明真实场景中的模型泛化能力。", styles["body"]))
    add_figure(story, figures[4], "图5 课程演示模型指标与特征重要性：指标来自 mock 数据，仅用于验证 Python 训练和看板链路。", styles)

    story.append(p("4 讨论", styles["h1"]))
    for para in [
        "第一，真实样本数不足是当前最大限制。17 条 completed 记录可以证明系统已能真实采集，但不适合划分稳定训练集和测试集。后续至少应补充到 50 条以上，并检查三类标签是否平衡。",
        "第二，手机运动特征存在兼容性限制。Web 页面无法保证锁屏或后台时持续采集传感器数据，且不同手机和浏览器支持程度不同。因此本项目将 motion_features 设计为可选输入，并在特征工程中显式记录 motion_available。",
        "第三，efficiency_score 是主观自评标签，不等同于客观成绩或真实学习成果。模型预测的是“自评学习效率等级”，适合课程实验和个人反思，不应推广为对所有大学生学习效率的普遍判断。",
        "第四，本项目允许使用 AI 工具辅助整理材料，但报告中尽量使用系统截图、流程图和数据表作为主图，避免使用与项目无关的生成式插画，降低明显 AI 痕迹。",
    ]:
        story.append(p(para, styles["body"]))

    story.append(p("5 总结", styles["h1"]))
    story.append(p("本课程设计完成了一个基于手机轻量感知与自报告数据的大学生学习效率预测系统原型。系统覆盖手机端学习记录、FastAPI 后端存储、训练数据导出、质量检查、随机森林三分类、预测接口和分析看板展示，体现了 Python 后端开发、数据处理、机器学习建模和物联网轻量感知的综合应用。当前版本已经适合课程答辩演示；后续可在增加真实样本、优化传感器采集稳定性、补充模型对比实验和完善部署环境后继续迭代。", styles["body"]))

    story.append(p("参考文献", styles["h1"]))
    refs = [
        "Python Software Foundation. Python Language Reference.",
        "FastAPI Documentation. FastAPI framework, high performance, easy to learn.",
        "Vue.js Documentation. The Progressive JavaScript Framework.",
        "scikit-learn developers. scikit-learn: Machine Learning in Python.",
        "Breiman, L. Random Forests. Machine Learning, 2001.",
        "MDN Web Docs. DeviceMotionEvent.",
        "SQLAlchemy Documentation. The Python SQL Toolkit and Object Relational Mapper.",
    ]
    for idx, item in enumerate(refs, 1):
        story.append(p(f"[{idx}] {item}", styles["ref"]))

    story.append(p("成员与分工", styles["h1"]))
    story.append(p(f"论文开头成员：{MEMBER_1}；{MEMBER_2}。", styles["body"]))
    division = [
        ["成员", "主要分工", "占比"],
        [MEMBER_1, "系统需求分析、后端接口、数据字典、模型训练脚本、论文结果整理", "50%"],
        [MEMBER_2, "前端页面、手机运动采集、看板展示、PPT 制作、演示与材料整理", "50%"],
    ]
    story.append(make_table(division, [0.28, 0.52, 0.20], styles))
    story.append(p("论文结尾成员：成员一姓名、成员二姓名。提交前请替换为真实姓名、学号和实际分工比例。", styles["body"]))

    doc.build(story, onFirstPage=page_header_footer, onLaterPages=page_header_footer)
    return out


def ppt_rgb(hex_value: str) -> RGBColor:
    return RGBColor(*hex_to_rgb(hex_value))


def add_slide_number(slide, idx: int) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(9.28), Inches(5.1), Inches(0.36), Inches(0.36))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ppt_rgb(ACCENT)
    shape.line.color.rgb = ppt_rgb(ACCENT)
    tf = shape.text_frame
    tf.clear()
    p0 = tf.paragraphs[0]
    run = p0.add_run()
    run.text = f"{idx:02d}"
    run.font.name = "Arial"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)
    p0.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size: int = 18, color: str = INK, bold: bool = False, align=PP_ALIGN.LEFT) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = ppt_rgb(color)


def add_title(slide, title: str, idx: int | None = None) -> None:
    add_text(slide, title, 0.55, 0.35, 8.3, 0.45, size=26, color=PRIMARY, bold=True)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(0.88), Inches(1.1), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = ppt_rgb(ACCENT)
    line.line.color.rgb = ppt_rgb(ACCENT)
    if idx:
        add_slide_number(slide, idx)


def add_card(slide, x, y, w, h, title, body, fill="#FFFFFF", border="#CBD5E1") -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ppt_rgb(fill)
    shape.line.color.rgb = ppt_rgb(border)
    add_text(slide, title, x + 0.18, y + 0.14, w - 0.36, 0.28, size=15, color=PRIMARY, bold=True)
    add_text(slide, body, x + 0.18, y + 0.52, w - 0.36, h - 0.62, size=11, color=INK)


def add_bullets(slide, items: list[str], x, y, w, h, size=14) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, item in enumerate(items):
        para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        para.text = item
        para.level = 0
        para.font.name = "Microsoft YaHei"
        para.font.size = Pt(size)
        para.font.color.rgb = ppt_rgb(INK)
        para.space_after = Pt(8)


def build_pptx(figures: list[Path]) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    blank = prs.slide_layouts[6]

    def bg(slide, color=BG):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = ppt_rgb(color)

    # 1 Cover
    slide = prs.slides.add_slide(blank)
    bg(slide, "#F8FAFC")
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(3.2), Inches(5.625)).fill.solid()
    slide.shapes[0].fill.fore_color.rgb = ppt_rgb(PRIMARY)
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(3.2), Inches(0), Inches(0.08), Inches(5.625)).fill.solid()
    slide.shapes[1].fill.fore_color.rgb = ppt_rgb(ACCENT)
    add_text(slide, "Python 期末课程设计", 0.45, 0.62, 2.25, 0.35, size=15, color="#FFFFFF")
    add_text(slide, TITLE, 3.65, 1.2, 5.8, 1.45, size=31, color=PRIMARY, bold=True)
    add_text(slide, "手机端学习打卡 · 自报告数据 · 随机森林三分类 · 分析看板", 3.7, 2.82, 5.5, 0.45, size=17, color=INK)
    add_text(slide, f"{CLASS_NAME}\\n{MEMBER_1}\\n{MEMBER_2}", 3.72, 3.78, 4.5, 0.9, size=13, color=MUTED)
    add_text(slide, "课程演示口径：真实数据质检 + mock 数据流程验证", 0.45, 4.62, 2.35, 0.5, size=12, color="#FFFFFF")

    # 2 Background
    slide = prs.slides.add_slide(blank); bg(slide); add_title(slide, "背景与课程设计目标", 2)
    add_card(slide, 0.65, 1.25, 2.55, 2.8, "问题", "只记录学习时长不足以解释学习效率差异。", blend(SECONDARY, 0.82), SECONDARY)
    add_card(slide, 3.55, 1.25, 2.55, 2.8, "方案", "用手机网页记录学习会话、自报告状态和可选运动特征。", "#FFFFFF", "#CBD5E1")
    add_card(slide, 6.45, 1.25, 2.55, 2.8, "目标", "完成可运行、可演示、可解释的 Python 课程设计原型。", blend(LIGHT, 0.62), "#D9A441")
    add_text(slide, "定位：不是成熟预测软件，而是课程级全栈 + 数据处理 + 机器学习闭环。", 0.8, 4.55, 8.4, 0.35, size=14, color=MUTED)

    # 3 Architecture
    slide = prs.slides.add_slide(blank); bg(slide); add_title(slide, "系统总体架构", 3)
    slide.shapes.add_picture(str(figures[0]), Inches(0.65), Inches(1.08), width=Inches(8.7))

    # 4 Collection
    slide = prs.slides.add_slide(blank); bg(slide); add_title(slide, "手机端采集流程", 4)
    slide.shapes.add_picture(str(figures[1]), Inches(0.58), Inches(1.05), width=Inches(8.8))

    # 5 Backend
    slide = prs.slides.add_slide(blank); bg(slide); add_title(slide, "后端 API 与数据存储", 5)
    modules = [
        ("用户", "simple-login\\n昵称进入"),
        ("学习记录", "start / end\\nlist / detail"),
        ("运动特征", "upload / get\\nupsert 保存"),
        ("模型", "train / predict\\nmetrics / importance"),
        ("分析", "overview / trend\\nfactor-analysis"),
    ]
    for i, (t, b) in enumerate(modules):
        add_card(slide, 0.55 + i * 1.82, 1.28, 1.55, 2.5, t, b, "#FFFFFF", SECONDARY if i % 2 == 0 else "#CBD5E1")
    add_text(slide, "数据库核心表：users、study_sessions、motion_features、predictions；删除历史记录另有备份表。", 0.8, 4.28, 8.2, 0.45, size=14, color=MUTED)

    # 6 Model
    slide = prs.slides.add_slide(blank); bg(slide); add_title(slide, "特征工程与模型方法", 6)
    slide.shapes.add_picture(str(figures[2]), Inches(0.58), Inches(1.03), width=Inches(8.85))

    # 7 Data
    slide = prs.slides.add_slide(blank); bg(slide); add_title(slide, "数据来源与实验口径", 7)
    add_card(slide, 0.75, 1.25, 3.9, 2.85, "真实数据", "completed = 17\\nlow=4, medium=11, high=2\\nmotion 缺失 11 条\\n用途：质检说明", "#FFFFFF", SECONDARY)
    add_card(slide, 5.3, 1.25, 3.9, 2.85, "mock 数据", "total = 120\\n三类标签各 40 条\\n用途：训练/预测/看板流程验证\\n不能写成真实规律", blend(ACCENT, 0.84), ACCENT)
    add_text(slide, "结论边界：少于 30 条真实 completed 记录时，不形成正式模型有效性结论。", 0.8, 4.45, 8.3, 0.35, size=14, color=PRIMARY, bold=True)

    # 8 UI and dashboard
    slide = prs.slides.add_slide(blank); bg(slide); add_title(slide, "系统界面与看板结果", 8)
    slide.shapes.add_picture(str(figures[3]), Inches(0.45), Inches(0.98), width=Inches(9.1))

    # 9 Model results
    slide = prs.slides.add_slide(blank); bg(slide); add_title(slide, "模型指标与特征重要性", 9)
    slide.shapes.add_picture(str(figures[4]), Inches(0.5), Inches(1.0), width=Inches(9.0))

    # 10 Demo and risks
    slide = prs.slides.add_slide(blank); bg(slide); add_title(slide, "演示流程与风险处理", 10)
    add_bullets(slide, ["1. 输入昵称进入系统", "2. 点击开始学习并计时", "3. 结束后填写自报告表单", "4. 查看历史记录与分析看板", "5. 生成预测结果和规则建议"], 0.8, 1.25, 3.9, 3.6, size=14)
    add_card(slide, 5.25, 1.18, 3.95, 1.02, "传感器兼容性", "运动特征可缺失，不阻断主流程。", blend(SECONDARY, 0.82), SECONDARY)
    add_card(slide, 5.25, 2.45, 3.95, 1.02, "数据量不足", "明确区分真实数据和 mock 数据。", blend(LIGHT, 0.62), "#D9A441")
    add_card(slide, 5.25, 3.72, 3.95, 1.02, "部署环境波动", "准备截图、样例数据和代码包。", blend(ACCENT, 0.84), ACCENT)

    # 11 Summary
    slide = prs.slides.add_slide(blank); bg(slide, "#F8FAFC"); add_title(slide, "总结与成员分工", 11)
    add_card(slide, 0.75, 1.2, 4.1, 2.95, "完成内容", "手机端采集、FastAPI 后端、数据导出、质量检查、随机森林演示、预测接口和看板展示。", "#FFFFFF", SECONDARY)
    add_card(slide, 5.25, 1.2, 4.1, 2.95, "分工占比", f"{MEMBER_1}：50%\\n{MEMBER_2}：50%\\n提交前请替换真实姓名学号。", "#FFFFFF", ACCENT)
    add_text(slide, "谢谢老师！", 0.8, 4.55, 8.3, 0.5, size=30, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)

    out = OUT / PPTX_NAME
    prs.save(out)
    return out


def package_code() -> Path:
    out = OUT / ZIP_NAME
    include_paths = [
        "README.md",
        "AGENTS.md",
        "docker-compose.yml",
        ".gitignore",
        "backend/app",
        "backend/tests",
        "backend/requirements.txt",
        "backend/README.md",
        "backend/.env.example",
        "frontend/src",
        "frontend/tests",
        "frontend/package.json",
        "frontend/package-lock.json",
        "frontend/vite.config.js",
        "frontend/index.html",
        "frontend/README.md",
        "frontend/.env.example",
        "ml",
        "docs/PLAN.md",
        "docs/TASKS.md",
        "docs/API_SPEC.md",
        "docs/DATA_DICTIONARY.md",
        "docs/DATA_PREP.md",
        "docs/EXPERIMENT_LOG.md",
        "docs/PAGE_FLOW.md",
        "docs/DEMO_SCRIPT.md",
        "docs/DATA_COLLECTION_PROTOCOL.md",
        "docs/SCREENSHOT_CHECKLIST.md",
        "docs/ai_outputs",
        "data/demo",
        "data/mock",
        "data/processed",
        "data/real",
        "models/latest.joblib",
    ]
    excluded_parts = {
        "__pycache__",
        ".pytest_cache",
        ".venv",
        "node_modules",
        "dist",
        ".git",
        ".git-local",
        ".agents",
        ".codex",
    }
    with zipfile.ZipFile(out, "w", compression=zipfile.ZIP_DEFLATED, compresslevel=9) as zf:
        for rel in include_paths:
            path = ROOT / rel
            if not path.exists():
                continue
            if path.is_file():
                zf.write(path, f"Study-efficiency/{rel}")
                continue
            for file in path.rglob("*"):
                if not file.is_file():
                    continue
                parts = set(file.relative_to(ROOT).parts)
                if parts & excluded_parts:
                    continue
                zf.write(file, f"Study-efficiency/{file.relative_to(ROOT).as_posix()}")
    return out


def write_submission_note(pdf: Path, pptx: Path, zip_path: Path, figures: list[Path]) -> Path:
    note = OUT / "提交前检查说明.txt"
    total = sum(p.stat().st_size for p in [pdf, pptx, zip_path])
    lines = [
        "Python期末课程设计交付物检查说明",
        "",
        f"论文PDF：{pdf.name}",
        f"PPT：{pptx.name}",
        f"代码压缩包：{zip_path.name}",
        f"三项合计大小：{total / 1024 / 1024:.2f} MB",
        "",
        "提交前必须替换：",
        "1. 文件名中的“提交人学号”和“提交人姓名”。",
        "2. 论文和PPT中的成员姓名、学号、分工比例占位符。",
        "3. 如实际分工不是50%/50%，请改成真实比例，合计100%。",
        "",
        "口径提醒：",
        "真实 completed 样本为17条，仅用于质检说明；120条mock数据仅用于课程演示流程验证，不代表真实学习效率规律。",
        "",
        "已生成主图：",
    ]
    lines += [f"- {p.name}" for p in figures]
    note.write_text("\n".join(lines), encoding="utf-8")
    return note


def main() -> None:
    ensure_dirs()
    figures = generate_figures()
    pdf = build_pdf(figures)
    pptx = build_pptx(figures)
    zip_path = package_code()
    note = write_submission_note(pdf, pptx, zip_path, figures)
    print("Generated:")
    for path in [pdf, pptx, zip_path, note, *figures]:
        print(f"- {path.relative_to(ROOT)} ({path.stat().st_size} bytes)")


if __name__ == "__main__":
    main()
